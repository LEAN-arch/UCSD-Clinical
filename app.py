import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import datetime
import io
import time
import logging

from prophet import Prophet
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from sklearn.linear_model import LogisticRegression
from sklearn.preprocessing import OneHotEncoder
from streamlit_option_menu import option_menu

# ======================================================================================
# SECTION 1: APP CONFIGURATION & STYLING
# ======================================================================================
st.set_page_config(
    page_title="MCC CTO QA Command Center",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
    .main .block-container { padding: 1rem 2rem 2rem; max-width: 1400px; }
    
    /* Base style for all metric displays */
    .stMetric { 
        background-color: #FAFAFA; 
        border: 1px solid #E0E0E0; 
        border-left: 5px solid #005A9C; 
        border-radius: 8px; 
        padding: 15px; 
    }

    /* Style for the explanation text that will appear BELOW the box */
    .kpi-explanation {
        font-size: 13px;
        color: #555555;
        padding-top: 10px;
        padding-left: 5px;
    }
</style>
""", unsafe_allow_html=True)

# ======================================================================================
# SECTION 1.5: SUPPRESS VERBOSE LOGGING
# ======================================================================================
logging.getLogger('cmdstanpy').setLevel(logging.WARNING)
logging.getLogger('prophet').setLevel(logging.WARNING)

# ======================================================================================
# SECTION 2: ENHANCED DATA SIMULATION
# ======================================================================================
@st.cache_data(ttl=900)
def generate_master_data():
    np.random.seed(42)
    num_trials = 60
    trial_ids = [f'MCC-{t}-{i:03d}' for t in ['IIT', 'IND', 'COG'] for i in range(1, 11)] + [f'INDUSTRY-{c}-{i:03d}' for c in ['PFE', 'BMY', 'MRK'] for i in range(1, 11)]
    
    portfolio_data = {
        'Trial_ID': trial_ids,
        'Trial_Type': np.random.choice(['Investigator-Initiated (IIT)', 'Industry-Sponsored', 'Cooperative Group'], num_trials, p=[0.4, 0.4, 0.2]),
        'Phase': np.random.choice(['I', 'I/II', 'II', 'III'], num_trials, p=[0.2, 0.3, 0.4, 0.1]),
        'Disease_Team': np.random.choice(['Leukemia', 'Lung', 'Breast', 'GI', 'GU', 'Melanoma'], num_trials),
        'Status': np.random.choice(['Enrolling', 'Follow-up', 'Closed to Accrual', 'Suspended'], num_trials, p=[0.6, 0.2, 0.15, 0.05]),
        'Subjects_Enrolled': np.random.randint(5, 100, num_trials),
        'PI_Experience_Level': np.random.choice(['Expert', 'Intermediate', 'New'], num_trials, p=[0.3, 0.5, 0.2]),
        'Is_First_In_Human': np.random.choice([True, False], num_trials, p=[0.1, 0.9]),
        'Num_Sites': np.random.choice([1, 2, 5], num_trials, p=[0.8, 0.15, 0.05]),
        'Screen_Fail_Rate': np.random.uniform(0.10, 0.50, num_trials),
        'Avg_Accrual_Per_Month': np.random.uniform(0.5, 5.0, num_trials),
        'Data_Query_Rate': np.random.uniform(0.05, 0.25, num_trials),
        'Last_Audit_Date': pd.to_datetime([datetime.date.today() - datetime.timedelta(days=np.random.randint(30, 730)) for _ in range(num_trials)]),
        'Days_to_Activation': np.random.randint(60, 200, num_trials)
    }
    portfolio_df = pd.DataFrame(portfolio_data)

    pis_by_team = {team: [f'Dr. {team_initials}{i}' for i in range(1,5)] for team, team_initials in zip(portfolio_df['Disease_Team'].unique(), ['L', 'U', 'B', 'GI', 'GU', 'M'])}
    portfolio_df['PI_Name'] = portfolio_df.apply(lambda row: np.random.choice(pis_by_team[row['Disease_Team']]), axis=1)

    finding_categories = ['Informed Consent Process', 'Source Data Verification', 'Investigational Product Accountability', 'Regulatory Binder Mgmt', 'AE/SAE Reporting', 'Protocol Adherence']
    findings_data = {'Finding_ID': [f'FIND-{i:04d}' for i in range(1, 251)], 'Trial_ID': np.random.choice(portfolio_df['Trial_ID'], 250), 'Category': np.random.choice(finding_categories, 250, p=[0.3, 0.2, 0.15, 0.15, 0.1, 0.1]), 'Risk_Level': np.random.choice(['Critical', 'Major', 'Minor'], 250, p=[0.05, 0.35, 0.6]), 'CAPA_Status': np.random.choice(['Open', 'Pending Verification', 'Closed-Effective', 'Overdue'], 250, p=[0.15, 0.1, 0.7, 0.05]), 'Finding_Date': pd.to_datetime([datetime.date(2022, 1, 1) + datetime.timedelta(days=int(d)) for d in np.random.randint(0, 900, 250)])}
    findings_df = pd.DataFrame(findings_data).merge(portfolio_df[['Trial_ID', 'Disease_Team', 'Trial_Type', 'PI_Name']], on='Trial_ID', how='left')
    findings_df['CAPA_Plan_Date'] = findings_df.apply(lambda r: r['Finding_Date'] + datetime.timedelta(days=np.random.randint(2, 7)) if r['CAPA_Status'] != 'Open' else pd.NaT, axis=1)
    findings_df['CAPA_Closure_Date'] = findings_df.apply(lambda r: r['CAPA_Plan_Date'] + datetime.timedelta(days=np.random.randint(10, 60)) if r['CAPA_Status'] == 'Closed-Effective' else pd.NaT, axis=1)

    major_finding_trials = findings_df[findings_df['Risk_Level'].isin(['Major', 'Critical'])]['Trial_ID'].unique()
    portfolio_df['Had_Major_Finding'] = portfolio_df['Trial_ID'].isin(major_finding_trials).astype(int)
    portfolio_df.loc[portfolio_df['Had_Major_Finding'] == 1, 'Screen_Fail_Rate'] *= 1.5
    portfolio_df.loc[portfolio_df['Had_Major_Finding'] == 1, 'Avg_Accrual_Per_Month'] *= 0.7
    portfolio_df.loc[portfolio_df['Had_Major_Finding'] == 1, 'Data_Query_Rate'] *= 1.8
    portfolio_df.loc[portfolio_df['Had_Major_Finding'] == 1, 'Days_to_Activation'] += 45

    auditors = ['Jane Doe, RN', 'John Smith, PhD', 'Maria Garcia, MPH', 'Kevin Lee, CCRC']
    team_data = {
        'Auditor': auditors,
        'Audits_Conducted_YTD': np.random.randint(15, 30, len(auditors)),
        'Avg_Report_Turnaround_Days': np.random.uniform(8, 20, len(auditors)),
        'GCP_Certification_Status': np.random.choice(['Current', 'Expires <90d'], len(auditors), p=[0.75, 0.25]),
        'IIT_Oversight_Skill': np.random.randint(3, 6, len(auditors)),
        'FDA_Inspection_Mgmt_Skill': np.random.randint(2, 5, len(auditors))
    }
    team_df = pd.DataFrame(team_data)

    initiatives_data = {'Initiative': ['eQMS Implementation', 'Auditor Training Program Revamp', 'Inspection Readiness Mock Audits', 'IIT Risk-Based Monitoring Plan'], 'Lead': ['Jane Doe, RN', 'John Smith, PhD', 'Maria Garcia, MPH', 'Kevin Lee, CCRC'], 'Status': ['On Track', 'At Risk', 'Completed', 'On Track'], 'Percent_Complete': [60, 85, 100, 30], 'Start_Date': pd.to_datetime(['2023-01-15', '2023-03-01', '2023-06-01', '2023-09-01']), 'End_Date': pd.to_datetime(['2024-06-30', '2023-11-30', '2023-08-31', '2024-03-31']), 'Budget_USD': [75000, 15000, 25000, 10000], 'Spent_USD': [40000, 14000, 23500, 2500]}
    initiatives_df = pd.DataFrame(initiatives_data)

    audits_data = []
    for index, finding in findings_df.iterrows():
        audits_data.append({
            'Audit_ID': f"AUDIT-{finding['Finding_ID']}", 'Trial_ID': finding['Trial_ID'],
            'Auditor': np.random.choice(team_df['Auditor']),
            'Audit_Date': finding['Finding_Date'] - datetime.timedelta(days=np.random.randint(1, 5)),
            'Turnaround_Time': np.random.uniform(7, 21)
        })
    audits_df = pd.DataFrame(audits_data)
    findings_df['Audit_ID'] = findings_df['Finding_ID'].apply(lambda x: f"AUDIT-{x}")
    findings_df['Is_Proactive'] = np.random.choice([True, False], len(findings_df), p=[0.3, 0.7])
    portfolio_df['Total_SAEs'] = np.random.randint(0, 15, len(portfolio_df))
    portfolio_df['Overdue_SAE_Reports'] = portfolio_df.apply(lambda r: np.random.randint(0, (r['Total_SAEs'] // 2) + 1) if r['Total_SAEs'] > 0 else 0, axis=1)
    
    return portfolio_df, findings_df, team_df, initiatives_df, audits_df

# ======================================================================================
# SECTION 3: ANALYTICAL & PLOTTING FUNCTIONS
# ======================================================================================
@st.cache_resource
def get_risk_model_with_importance(_portfolio_df):
    features, target = ['Trial_Type', 'Phase', 'PI_Experience_Level', 'Is_First_In_Human', 'Num_Sites'], 'Had_Major_Finding'
    X, y = _portfolio_df[features], _portfolio_df[target]
    categorical_features = ['Trial_Type', 'Phase', 'PI_Experience_Level']
    encoder = OneHotEncoder(handle_unknown='ignore', sparse_output=False)
    X_encoded = pd.DataFrame(encoder.fit_transform(X[categorical_features]), columns=encoder.get_feature_names_out(categorical_features))
    X_final = pd.concat([X.drop(columns=categorical_features).reset_index(drop=True), X_encoded], axis=1)
    model = LogisticRegression(max_iter=1000, class_weight='balanced', random_state=42)
    model.fit(X_final, y)
    importance = pd.DataFrame(data=model.coef_[0], index=X_final.columns, columns=['Coefficient']).sort_values(by='Coefficient', ascending=False)
    return model, encoder, X_final.columns, importance

@st.cache_resource
def get_finding_category_model(_portfolio_df, _findings_df):
    major_findings = _findings_df[_findings_df['Risk_Level'].isin(['Major', 'Critical'])]
    if major_findings.empty: return None, None, None
    first_major = major_findings.sort_values('Finding_Date').groupby('Trial_ID').first().reset_index()
    model_data = pd.merge(_portfolio_df, first_major[['Trial_ID', 'Category']], on='Trial_ID', how='left')
    model_data['Category'].fillna('No Major Finding', inplace=True)
    features, target = ['Trial_Type', 'Phase', 'PI_Experience_Level', 'Is_First_In_Human', 'Num_Sites'], 'Category'
    X, y = model_data[features], model_data[target]
    categorical_features = ['Trial_Type', 'Phase', 'PI_Experience_Level']
    encoder = OneHotEncoder(handle_unknown='ignore', sparse_output=False)
    X_encoded = pd.DataFrame(encoder.fit_transform(X[categorical_features]), columns=encoder.get_feature_names_out(categorical_features))
    X_final = pd.concat([X.drop(columns=categorical_features).reset_index(drop=True), X_encoded], axis=1)
    model = LogisticRegression(max_iter=1000, class_weight='balanced', multi_class='multinomial', random_state=42)
    model.fit(X_final, y)
    return model, encoder, X_final.columns

@st.cache_data(ttl=3600)
def generate_prophet_forecast(_findings_df):
    df_prophet = _findings_df[['Finding_Date']].copy()
    df_prophet = df_prophet.rename(columns={'Finding_Date': 'ds'})
    monthly_df = df_prophet.set_index('ds').resample('ME').size().reset_index()
    monthly_df = monthly_df.rename(columns={0: 'y'})
    if monthly_df.empty:
        return pd.DataFrame(), monthly_df
    model = Prophet(yearly_seasonality=True, daily_seasonality=False)
    model.fit(monthly_df)
    future = model.make_future_dataframe(periods=12, freq='ME')
    forecast = model.predict(future)
    return forecast, monthly_df

@st.cache_data
def analyze_capa_effectiveness(findings_df):
    closed_capas = findings_df[findings_df['CAPA_Status'] == 'Closed-Effective'].copy()
    if closed_capas.empty: return pd.DataFrame(columns=['Category', 'Recurrence_Rate'])
    recurrence_data = []
    for index, capa in closed_capas.iterrows():
        subsequent_findings = findings_df[(findings_df['Trial_ID'] == capa['Trial_ID']) & (findings_df['Category'] == capa['Category']) & (findings_df['Finding_Date'] > capa['Finding_Date'])]
        recurrence_data.append({'Category': capa['Category'], 'Recurrence': int(not subsequent_findings.empty)})
    if not recurrence_data: return pd.DataFrame(columns=['Category', 'Recurrence_Rate'])
    recurrence_df = pd.DataFrame(recurrence_data)
    recurrence_rate = recurrence_df.groupby('Category')['Recurrence'].mean().reset_index()
    recurrence_rate.rename(columns={'Recurrence': 'Recurrence_Rate'}, inplace=True)
    return recurrence_rate.sort_values('Recurrence_Rate', ascending=False)

def plot_prophet_forecast_sme(forecast, monthly_df):
    fig = go.Figure()
    if not monthly_df.empty:
        fig.add_trace(go.Scatter(x=monthly_df['ds'], y=monthly_df.iloc[:, 1], mode='markers', name='Actual Findings', marker=dict(color='#005A9C', size=8), hovertemplate="<b>%{x|%B %Y}</b><br>Actual Findings: %{y}<extra></extra>"))
    if not forecast.empty:
        fig.add_trace(go.Scatter(x=forecast['ds'], y=forecast['yhat'], mode='lines', name='Forecast', line=dict(color='#3EC1D3', dash='dash'), hovertemplate="<b>%{x|%B %Y}</b><br>Forecasted: %{y:.1f}<extra></extra>"))
        fig.add_trace(go.Scatter(x=forecast['ds'], y=forecast['yhat_upper'], fill=None, mode='lines', line_color='rgba(62,193,211,0.2)', showlegend=False))
        fig.add_trace(go.Scatter(x=forecast['ds'], y=forecast['yhat_lower'], fill='tonexty', mode='lines', line_color='rgba(62,193,211,0.2)', name='Uncertainty', hoverinfo='none'))
        fig.add_trace(go.Scatter(x=forecast['ds'], y=forecast['trend'], mode='lines', name='Overall Trend', line=dict(color='#FFC72C', width=3), hovertemplate="Overall Trend: %{y:.1f}<extra></extra>"))
    
    if not monthly_df.empty:
        last_actual_date = monthly_df['ds'].iloc[-1]
        fig.add_vline(x=last_actual_date, line_width=1, line_dash="dot", line_color="grey")
        fig.add_annotation(
            x=last_actual_date, y=1, yref="paper",
            text="Last Actual", showarrow=False,
            xanchor="left", yanchor="bottom", xshift=5
        )
    fig.update_layout(title='<b>12-Month Forecast of Audit Findings with Trend Analysis</b>', xaxis_title=None, yaxis_title='Number of Findings', plot_bgcolor='white', legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
    return fig

def plot_spc_chart_sme(df, date_col, category_col, value, title):
    df_filtered = df[df[category_col] == value].copy()
    if df_filtered.empty: return go.Figure().update_layout(title=f'<b>{title}</b><br>No data available.')
    df_filtered[date_col] = pd.to_datetime(df_filtered[date_col])
    df_filtered = df_filtered.set_index(date_col).sort_index()
    monthly_counts = df_filtered.resample('ME').size().reset_index(name='findings')
    monthly_counts['month'] = monthly_counts[date_col].dt.to_period('M').astype(str)
    if monthly_counts.empty or monthly_counts['findings'].sum() == 0: return go.Figure().update_layout(title=f'<b>{title}</b><br>No data available.')
    p_bar, std_dev = monthly_counts['findings'].mean(), np.sqrt(monthly_counts['findings'].mean())
    UCL, LCL = p_bar + 3 * std_dev, max(0, p_bar - 3 * std_dev)
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=monthly_counts['month'], y=[UCL]*len(monthly_counts), mode='lines', line=dict(color='rgba(255, 100, 100, 0.5)'), showlegend=False, hoverinfo='none'))
    fig.add_trace(go.Scatter(x=monthly_counts['month'], y=[LCL]*len(monthly_counts), mode='lines', line=dict(color='rgba(255, 100, 100, 0.5)'), fill='tonexty', fillcolor='rgba(0, 176, 246, 0.1)', name='Common Cause Variation Zone', hoverinfo='none'))
    fig.add_trace(go.Scatter(x=monthly_counts['month'], y=[p_bar]*len(monthly_counts), mode='lines', name='Process Mean', line=dict(color='green', dash='dot')))
    fig.add_trace(go.Scatter(x=monthly_counts['month'], y=[UCL]*len(monthly_counts), mode='lines', name='Upper Control Limit', line=dict(color='red', dash='dash')))
    fig.add_trace(go.Scatter(x=monthly_counts['month'], y=[LCL]*len(monthly_counts), mode='lines', name='Lower Control Limit', line=dict(color='red', dash='dash')))
    fig.add_trace(go.Scatter(x=monthly_counts['month'], y=monthly_counts['findings'], mode='lines+markers', name='Monthly Findings', line=dict(color='#005A9C'), hovertemplate="<b>%{x}</b><br>Findings: %{y}<br>Status: In Control<extra></extra>"))
    out_of_control = monthly_counts[monthly_counts['findings'] > UCL]
    if not out_of_control.empty:
        fig.add_trace(go.Scatter(x=out_of_control['month'], y=out_of_control['findings'], mode='markers', name='Out of Control Signal', marker=dict(color='red', size=12, symbol='x', line=dict(width=3)), hovertemplate="<b>%{x}</b><br>Findings: %{y}<br>Status: <b>Out of Control</b><extra></extra>"))
    fig.update_layout(title=f'<b>{title}</b>', xaxis_title=None, yaxis_title='Number of Findings', plot_bgcolor='white', height=450, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
    return fig

def generate_ppt_report(kpi_data, spc_fig, findings_table_df):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "MCC CTO Quality Assurance Executive Summary"
    slide.placeholders[1].text = f"Report Generated: {datetime.date.today().strftime('%Y-%m-%d')}"
    kpi_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(kpi_slide_layout)
    slide.shapes.title.text = "QA Program Health Dashboard"
    positions = [(Inches(1), Inches(1.5)), (Inches(5), Inches(1.5)), (Inches(9), Inches(1.5)), (Inches(13), Inches(1.5))]
    for i, (kpi_title, kpi_val, kpi_delta) in enumerate(kpi_data):
        txBox = slide.shapes.add_textbox(positions[i][0], positions[i][1], Inches(3.5), Inches(2))
        tf = txBox.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.text = kpi_title; p1.font.bold = True; p1.font.size = Pt(20)
        p2 = tf.add_paragraph(); p2.text = str(kpi_val); p2.font.size = Pt(44); p2.font.bold = True
        p3 = tf.add_paragraph(); p3.text = str(kpi_delta); p3.font.size = Pt(16)
    chart_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(chart_slide_layout)
    slide.shapes.title.text = "Systemic Process Control (SPC) Analysis"
    try:
        image_stream = io.BytesIO()
        spc_fig.write_image(image_stream, format='png', scale=2)
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(14))
    except (RuntimeError, ValueError) as e:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(4))
        tf = txBox.text_frame
        p = tf.add_paragraph(); p.text = "Chart Generation Failed: Missing Dependency"; p.font.bold = True; p.font.size = Pt(24); p.font.color.rgb = RGBColor(255, 0, 0)
        p = tf.add_paragraph(); p.text = ("The 'Kaleido' library is required to export this chart. If running locally, install it (`pip install kaleido`). The deployed app environment may be missing this dependency."); p.font.size = Pt(18)
    table_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(table_slide_layout)
    slide.shapes.title.text = "High-Priority Open Findings (Critical/Major)"
    rows, cols = findings_table_df.shape[0] + 1, findings_table_df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(14), Inches(0.5) * rows).table
    for col_idx, col_name in enumerate(findings_table_df.columns): table.cell(0, col_idx).text = col_name
    for ppt_row_idx, df_row in enumerate(findings_table_df.itertuples(index=False), start=1):
        for col_idx, cell_data in enumerate(df_row): table.cell(ppt_row_idx, col_idx).text = str(cell_data)
    ppt_stream = io.BytesIO(); prs.save(ppt_stream); ppt_stream.seek(0)
    return ppt_stream

def plot_pi_findings_barchart_sme(pi_findings, pi_name):
    if pi_findings.empty: return go.Figure().update_layout(title=f'No findings recorded for {pi_name}', annotations=[dict(text="No Data", x=0.5, y=0.5, showarrow=False, font_size=20)])
    category_counts = pi_findings['Category'].value_counts().reset_index()
    category_counts.columns = ['Category', 'Count']
    fig = px.bar(category_counts.sort_values('Count'), x='Count', y='Category', orientation='h', title=f'<b>Finding Breakdown for {pi_name}</b>', text='Count', labels={'Count': 'Number of Findings', 'Category': 'Finding Category'})
    fig.update_traces(marker_color='#005A9C', textposition='outside')
    fig.update_layout(plot_bgcolor='white', yaxis={'categoryorder':'total ascending'})
    return fig

def plot_auditor_strain_sme(team_df):
    team_df['Skill_Factor'] = team_df['IIT_Oversight_Skill'] + team_df['FDA_Inspection_Mgmt_Skill']
    team_df['Strain'] = (team_df['Audits_Conducted_YTD'] * team_df['Avg_Report_Turnaround_Days']) / (team_df['Skill_Factor'] + 1)
    fig = px.scatter(team_df, x='Audits_Conducted_YTD', y='Avg_Report_Turnaround_Days', size='Strain', color='Strain', text='Auditor', title='<b>Auditor Performance & Workload Quadrant Analysis</b>', labels={'Audits_Conducted_YTD': 'Audits Conducted (Workload)', 'Avg_Report_Turnaround_Days': 'Avg. Report Turnaround Time (Efficiency)'}, color_continuous_scale=px.colors.sequential.OrRd)
    mean_x, mean_y = team_df['Audits_Conducted_YTD'].mean(), team_df['Avg_Report_Turnaround_Days'].mean()
    fig.add_hline(y=mean_y, line_dash="dot", line_color="grey", annotation_text="Avg. Efficiency")
    fig.add_vline(x=mean_x, line_dash="dot", line_color="grey", annotation_text="Avg. Workload")
    fig.add_annotation(x=mean_x*1.2, y=mean_y*1.2, text="<b>High Strain</b><br>(At Risk of Burnout)", showarrow=False, font=dict(color="firebrick"))
    fig.add_annotation(x=mean_x*1.2, y=mean_y*0.8, text="<b>Top Performers</b><br>(Potential Mentors)", showarrow=False, font=dict(color="darkgreen"))
    fig.add_annotation(x=mean_x*0.8, y=mean_y*0.8, text="<b>Underutilized</b><br>(Capacity for Growth)", showarrow=False, font=dict(color="darkblue"))
    fig.add_annotation(x=mean_x*0.8, y=mean_y*1.2, text="<b>Needs Coaching</b><br>(Efficiency Opportunity)", showarrow=False, font=dict(color="goldenrod"))
    fig.update_traces(textposition='top center', hovertemplate="<b>%{text}</b><br><br>Workload (Audits YTD): %{x}<br>Efficiency (Avg Turnaround): %{y:.1f} days<br><b>Skill-Weighted Strain Index: %{marker.size:.2f}</b><extra></extra>")
    fig.update_layout(plot_bgcolor='white', height=500)
    return fig

# ======================================================================================
# SECTION 4: UI PAGE RENDERING FUNCTIONS
# ======================================================================================
def render_command_center(portfolio_df, findings_df, team_df):
    st.subheader("Executive Command Center", divider="blue")
    st.markdown("A strategic overview of the QA program's current status, efficiency, and highest priority items.")

    risk_weights = {'Critical': 10, 'Major': 5, 'Minor': 1}
    findings_df['Risk_Score'] = findings_df['Risk_Level'].map(risk_weights)
    open_findings = findings_df[~findings_df['CAPA_Status'].isin(['Closed-Effective'])].copy()

    # --- Program Health & Risk ---
    st.markdown("##### Program Health & Risk")
    kpi_col1, kpi_col2, kpi_col3 = st.columns(3)

    with kpi_col1:
        total_risk_score = int(open_findings['Risk_Score'].sum())
        st.metric("Portfolio-wide Risk Score", f"{total_risk_score}", f"{open_findings[open_findings['Risk_Level'] == 'Critical'].shape[0]} Open Criticals", "inverse")
        st.markdown("<p class='kpi-explanation'>A weighted sum of all open findings. A higher score indicates greater overall portfolio risk.</p>", unsafe_allow_html=True)

    with kpi_col2:
        overdue_major_capas = findings_df[(findings_df['CAPA_Status'] == 'Overdue') & (findings_df['Risk_Level'] != 'Minor')].shape[0]
        readiness_score = max(0, 100 - (overdue_major_capas * 10) - (open_findings[open_findings['Risk_Level'] == 'Critical'].shape[0] * 5))
        st.metric("Inspection Readiness Index", f"{readiness_score}%", f"{overdue_major_capas} Overdue Major CAPAs", "inverse")
        st.markdown("<p class='kpi-explanation'>A 0-100% score based on overdue and critical findings. Represents preparedness for a regulatory inspection.</p>", unsafe_allow_html=True)

    with kpi_col3:
        if 'Strain' not in team_df.columns:
            team_df['Skill_Factor'] = team_df['IIT_Oversight_Skill'] + team_df['FDA_Inspection_Mgmt_Skill']
            team_df['Strain'] = (team_df['Audits_Conducted_YTD'] * team_df['Avg_Report_Turnaround_Days']) / (team_df['Skill_Factor'] + 1)
        avg_strain = team_df['Strain'].mean()
        st.metric("Avg. Resource Strain Index", f"{avg_strain:.2f}", "Target < 2.5", "normal")
        st.markdown("<p class='kpi-explanation'>A metric combining auditor workload and efficiency. A higher value may indicate a risk of burnout or bottlenecks.</p>", unsafe_allow_html=True)

    # --- Program Velocity & Maturity ---
    st.markdown("##### Program Velocity & Maturity")
    kpi_col4, kpi_col5, kpi_col6, kpi_col7, kpi_col8, kpi_col9 = st.columns(6)

    with kpi_col4:
        open_findings['Finding_Date'] = pd.to_datetime(open_findings['Finding_Date'])
        open_findings['Age'] = (datetime.datetime.now() - open_findings['Finding_Date']).dt.days
        avg_capa_age = open_findings['Age'].mean()
        st.metric("Avg. Open CAPA Age", f"{avg_capa_age:.1f} days", "Target < 30", "inverse")

    with kpi_col5:
        total_saes = portfolio_df['Total_SAEs'].sum()
        overdue_saes = portfolio_df['Overdue_SAE_Reports'].sum()
        overdue_sae_rate = (overdue_saes / total_saes) * 100 if total_saes > 0 else 0
        st.metric("Overdue SAE Reporting", f"{overdue_sae_rate:.2f}%", f"{overdue_saes} Overdue", "inverse")

    with kpi_col6:
        enrolled_subjects = portfolio_df['Subjects_Enrolled'].sum()
        first_pass_rate = 1 - (len(findings_df) / (enrolled_subjects * 5)) if enrolled_subjects > 0 else 1.0
        st.metric("First Pass Quality", f"{first_pass_rate:.1%}", "Target > 95%", "normal")

    with kpi_col7:
        avg_query_rate = portfolio_df['Data_Query_Rate'].mean()
        data_integrity_score = (1 - avg_query_rate) * 100
        st.metric("Data Integrity Score", f"{data_integrity_score:.1f}%", f"{avg_query_rate:.2f} Q/Subj", "normal")

    with kpi_col8:
        proactive_audits = findings_df['Is_Proactive'].sum()
        total_audits = len(findings_df)
        proactive_ratio = (proactive_audits / total_audits) * 100 if total_audits > 0 else 0
        st.metric("Proactive Audit Ratio", f"{proactive_ratio:.1f}%", f"{proactive_audits} Proactive", "normal")

    with kpi_col9:
        current_certs = len(team_df[team_df['GCP_Certification_Status'] == 'Current'])
        total_auditors = len(team_df)
        team_readiness = (current_certs / total_auditors) * 100 if total_auditors > 0 else 0
        expiring_soon = total_auditors - current_certs
        st.metric("Team Readiness (GCP)", f"{team_readiness:.0f}%", f"{expiring_soon} Expiring", "off")
        
    st.markdown("<hr>", unsafe_allow_html=True)

    # --- Plotting Tabs ---
    plot_tabs = st.tabs(["🔥 Priority Alerts", "📊 Finding Backlog", "🧑‍🔬 Auditor Skills", "🗺️ Risk Treemap"])
    
    with plot_tabs[0]:
        st.markdown("##### High-Priority Alerts & Portfolio Status")
        col1, col2 = st.columns([1, 2])
        with col1:
            overdue_findings = findings_df[(findings_df['CAPA_Status'] == 'Overdue') & (findings_df['Risk_Level'] != 'Minor')].sort_values(by='Finding_Date').reset_index()
            if not overdue_findings.empty: st.error(f"**Overdue CAPA:** Finding `{overdue_findings.iloc[0]['Finding_ID']}` on trial `{overdue_findings.iloc[0]['Trial_ID']}` ({overdue_findings.iloc[0]['Risk_Level']}) is overdue.", icon="🔥")
            if 'Strain' not in team_df.columns:
                team_df['Skill_Factor'] = team_df['IIT_Oversight_Skill'] + team_df['FDA_Inspection_Mgmt_Skill']
                team_df['Strain'] = (team_df['Audits_Conducted_YTD'] * team_df['Avg_Report_Turnaround_Days']) / (team_df['Skill_Factor'] + 1)
            most_strained = team_df.sort_values(by='Strain', ascending=False).iloc[0]
            if most_strained['Strain'] > 3.0: st.warning(f"**Resource At Risk:** `{most_strained['Auditor']}` has a high Strain Index of `{most_strained['Strain']:.2f}`.", icon="⚠️")
            criticals_per_trial = findings_df[findings_df['Risk_Level'] == 'Critical'].groupby('Trial_ID').size().sort_values(ascending=False)
            if not criticals_per_trial.empty: st.error(f"**High-Risk Trial:** Trial `{criticals_per_trial.index[0]}` has `{criticals_per_trial.iloc[0]}` open critical findings.", icon="🔬")
        with col2:
            status_counts = portfolio_df['Status'].value_counts()
            fig = px.pie(status_counts, values=status_counts.values, names=status_counts.index, title='Active Clinical Trial Portfolio', hole=0.4, color_discrete_map={'Enrolling':'#005A9C', 'Follow-up':'#3EC1D3', 'Closed to Accrual':'#FFC72C', 'Suspended':'#E63946'})
            st.plotly_chart(fig, use_container_width=True)
            
    with plot_tabs[1]:
        st.markdown("##### Monthly Finding Velocity & Active Backlog Analysis")
        st.info("💡 **Expert Tip:** Compare the stacked 'Opened' bar to the solid 'Closed' bar each month. If 'Opened' is consistently taller, the backlog (blue line) will rise. Focus on months where the red (Critical) or orange (Major) segments are large.", icon="❓")
        findings_df['Finding_Month'] = findings_df['Finding_Date'].dt.to_period('M')
        opened_by_risk = pd.crosstab(findings_df['Finding_Month'], findings_df['Risk_Level'])
        for risk in ['Critical', 'Major', 'Minor']:
            if risk not in opened_by_risk.columns:
                opened_by_risk[risk] = 0
        opened_by_risk = opened_by_risk[['Critical', 'Major', 'Minor']]
        closed_by_month = findings_df.dropna(subset=['CAPA_Closure_Date'])
        if not closed_by_month.empty:
            closed_by_month['Closure_Month'] = closed_by_month['CAPA_Closure_Date'].dt.to_period('M')
            closed_counts = closed_by_month.groupby('Closure_Month').size().rename('Closed')
        else:
            closed_counts = pd.Series(name='Closed', dtype='int64')
        velocity_df = pd.concat([opened_by_risk, closed_counts], axis=1).fillna(0)
        velocity_df['Total_Opened'] = velocity_df[['Critical', 'Major', 'Minor']].sum(axis=1)
        velocity_df['Cumulative_Opened'] = velocity_df['Total_Opened'].cumsum()
        velocity_df['Cumulative_Closed'] = velocity_df['Closed'].cumsum()
        velocity_df['Active_Backlog'] = velocity_df['Cumulative_Opened'] - velocity_df['Cumulative_Closed']
        velocity_df.index = velocity_df.index.to_timestamp()
        fig = go.Figure()
        fig.add_trace(go.Bar(x=velocity_df.index, y=velocity_df['Critical'], name='Opened: Critical', marker_color='#E63946'))
        fig.add_trace(go.Bar(x=velocity_df.index, y=velocity_df['Major'], name='Opened: Major', marker_color='#FFC72C'))
        fig.add_trace(go.Bar(x=velocity_df.index, y=velocity_df['Minor'], name='Opened: Minor', marker_color='#A8DADC'))
        fig.add_trace(go.Bar(x=velocity_df.index, y=velocity_df['Closed'], name='Closed (All Risks)', marker_color='#457B9D', hovertemplate="<b>%{x|%B %Y}</b><br>Closed: %{y}<extra></extra>"))
        fig.add_trace(go.Scatter(x=velocity_df.index, y=velocity_df['Active_Backlog'], name='Active Backlog Size', mode='lines+markers', line=dict(color='#005A9C', width=3), yaxis='y2', hovertemplate="<b>%{x|%B %Y}</b><br>Backlog: %{y}<extra></extra>"))
        fig.update_layout(barmode='stack', title_text='<b>Monthly Finding Velocity vs. Active Backlog Trend</b>', yaxis=dict(title='Monthly Count'), yaxis2=dict(title='Total Open Backlog', overlaying='y', side='right', showgrid=False), legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1), plot_bgcolor='white')
        st.plotly_chart(fig, use_container_width=True)

    with plot_tabs[2]:
        st.markdown("##### Interactive Auditor Skill & Performance Matrix")
        skill_cols = [col for col in team_df.columns if '_Skill' in col]
        team_df['Overall_Skill'] = team_df[skill_cols].mean(axis=1)
        skill_options = {col.replace('_', ' ').replace(' Skill', ''): col for col in skill_cols}
        col1, col2 = st.columns([1, 2])
        with col1:
            st.info("💡 **Expert Tip:** Select a skill from the dropdown to color-code the auditors. This helps you find the right expert and assess their current workload and efficiency simultaneously.", icon="❓")
            selected_skill_label = st.selectbox("Color-code auditors by skill:", options=list(skill_options.keys()), index=0)
            selected_skill_col = skill_options[selected_skill_label]
        with col2:
            fig = px.scatter(
                team_df,
                x='Audits_Conducted_YTD',
                y='Avg_Report_Turnaround_Days',
                size='Overall_Skill',
                color=selected_skill_col,
                text='Auditor',
                title=f'<b>Auditor Performance (Colored by {selected_skill_label} Skill)</b>',
                labels={
                    'Audits_Conducted_YTD': 'Workload (Audits YTD)',
                    'Avg_Report_Turnaround_Days': 'Efficiency (Avg. Report Turnaround)',
                    selected_skill_col: f'{selected_skill_label} Level'
                },
                color_continuous_scale=px.colors.sequential.Blues,
                size_max=40,
                height=500
            )
            mean_x = team_df['Audits_Conducted_YTD'].mean()
            mean_y = team_df['Avg_Report_Turnaround_Days'].mean()
            fig.add_hline(y=mean_y, line_dash="dot", line_color="grey")
            fig.add_vline(x=mean_x, line_dash="dot", line_color="grey")
            fig.add_annotation(x=mean_x*1.2, y=mean_y*1.2, text="<b>Needs Coaching</b><br>(High Workload, Low Efficiency)", showarrow=False, font=dict(color="goldenrod"))
            fig.add_annotation(x=mean_x*1.2, y=mean_y*0.8, text="<b>Top Performers</b><br>(High Workload, High Efficiency)", showarrow=False, font=dict(color="darkgreen"))
            fig.add_annotation(x=mean_x*0.8, y=mean_y*0.8, text="<b>Growth Potential</b><br>(Low Workload, High Efficiency)", showarrow=False, font=dict(color="darkblue"))
            fig.add_annotation(x=mean_x*0.8, y=mean_y*1.2, text="<b>Efficiency Opportunity</b><br>(Low Workload, Low Efficiency)", showarrow=False, font=dict(color="firebrick"))
            fig.update_traces(
                textposition='top center',
                hovertemplate=(
                    "<b>%{text}</b><br><br>"
                    "Workload (Audits): %{x}<br>"
                    "Efficiency (Turnaround): %{y:.1f} days<br>"
                    "Overall Skill Score: %{marker.size:.1f}<br>"
                    f"{selected_skill_label} Skill: %{{customdata[0]}}<extra></extra>"
                ),
                customdata=team_df[[selected_skill_col]]
            )
            fig.update_layout(plot_bgcolor='white', coloraxis_colorbar=dict(title=f'{selected_skill_label}<br>Level'))
            st.plotly_chart(fig, use_container_width=True)
               
    with plot_tabs[3]:
        st.markdown("##### Interactive Risk Exploration Tool")
        risk_summary = findings_df.groupby('Trial_ID').agg(Risk_Score=('Risk_Score', 'sum'), Total_Findings=('Finding_ID', 'count')).reset_index()
        risk_counts = pd.crosstab(findings_df['Trial_ID'], findings_df['Risk_Level'])
        for risk in ['Critical', 'Major', 'Minor']:
            if risk not in risk_counts.columns:
                risk_counts[risk] = 0
        risk_counts = risk_counts.reset_index().rename(columns={'Critical': 'Critical_Count', 'Major': 'Major_Count', 'Minor': 'Minor_Count'})
        open_capa_counts = findings_df[findings_df['CAPA_Status'] != 'Closed-Effective'].groupby('Trial_ID').size().reset_index(name='Open_CAPA_Count')
        risk_model, encoder, model_features, _ = get_risk_model_with_importance(portfolio_df)
        X_all_cat = portfolio_df[['Trial_Type', 'Phase', 'PI_Experience_Level']]
        X_all_num = portfolio_df[['Is_First_In_Human', 'Num_Sites']]
        X_all_encoded = pd.DataFrame(encoder.transform(X_all_cat), columns=encoder.get_feature_names_out(X_all_cat.columns))
        X_all_final = pd.concat([X_all_num.reset_index(drop=True), X_all_encoded], axis=1).reindex(columns=model_features, fill_value=0)
        portfolio_df['Predicted_Risk'] = risk_model.predict_proba(X_all_final)[:, 1]
        risk_map_df = portfolio_df.merge(risk_summary, on='Trial_ID', how='left')
        risk_map_df = risk_map_df.merge(risk_counts, on='Trial_ID', how='left')
        risk_map_df = risk_map_df.merge(open_capa_counts, on='Trial_ID', how='left')
        risk_map_df.fillna(0, inplace=True)
        risk_map_df['Avg_Risk_Severity'] = risk_map_df.apply(lambda r: r['Risk_Score'] / r['Total_Findings'] if r['Total_Findings'] > 0 else 0, axis=1)

        col1, col2 = st.columns([1, 4])
        with col1:
            st.info("💡 **Expert Tip:** Use these controls to analyze risk from different perspectives. Switch to 'Predicted Risk' to be proactive!")
            view_selection = st.selectbox("Select Risk View (Box Size):", options=["Observed Risk (from past findings)", "Predicted Risk (from model)"], index=0)
            color_selection = st.selectbox("Select Metric (Box Color):", options=["Number of Critical Findings", "Number of Open CAPAs", "Average Finding Severity"], index=0)
        with col2:
            view_column_map = {"Observed Risk (from past findings)": "Risk_Score", "Predicted Risk (from model)": "Predicted_Risk"}
            color_column_map = {"Number of Critical Findings": "Critical_Count", "Number of Open CAPAs": "Open_CAPA_Count", "Average Finding Severity": "Avg_Risk_Severity"}
            view_col = view_column_map[view_selection]
            color_col = color_column_map[color_selection]
            fig = px.treemap(
                risk_map_df,
                path=[px.Constant("All Trials"), 'Disease_Team', 'Trial_Type', 'PI_Name'],
                values=view_col, color=color_col,
                custom_data=['Trial_ID', 'Risk_Score', 'Predicted_Risk', 'Critical_Count', 'Major_Count', 'Minor_Count', 'Open_CAPA_Count'],
                color_continuous_scale='Reds',
                title=f'<b>Portfolio Risk Analysis</b><br><i>Box Size by {view_selection} | Box Color by {color_selection}</i>'
            )
            fig.update_traces(hovertemplate=("<b>%{label}</b><br><br>" + "<b>Predicted Risk Score:</b> %{customdata[2]:.1%}<br>" + "<b>Observed Risk Score:</b> %{customdata[1]:.0f}<br>" + "<b>Critical Findings:</b> %{customdata[3]}<br>" + "<b>Major Findings:</b> %{customdata[4]}<br>" + "<b>Open CAPAs:</b> %{customdata[6]:.0f}" + "<extra></extra>"))
            fig.update_layout(margin=dict(t=80, l=25, r=25, b=25))
            st.plotly_chart(fig, use_container_width=True)

def render_predictive_analytics(findings_df, portfolio_df):
    st.subheader("Predictive Analytics & Forecasting", divider="blue")
    st.markdown("_This section utilizes predictive modeling to forecast future states and quantify inherent risk, enabling a proactive, data-driven approach to quality management._")
    
    forecast_data, actual_monthly_data = generate_prophet_forecast(findings_df)
    risk_model, encoder, model_features, importance_df = get_risk_model_with_importance(portfolio_df)

    with st.container(border=True):
        st.markdown("##### Time-Series Forecast of Audit Finding Volume")
        st.info("💡 **Expert Tip:** Is the overall yellow trend line increasing, decreasing, or flat? An increasing trend suggests systemic issues may be worsening, requiring strategic intervention beyond addressing individual findings.", icon="❓")
        st.plotly_chart(plot_prophet_forecast_sme(forecast_data, actual_monthly_data), use_container_width=True)

    with st.container(border=True):
        st.markdown("##### Inherent Risk Prediction for New Trials")
        st.info("💡 **Expert Tip:** Use this tool to triage new protocols. A predicted risk score > 60% may warrant assigning a more senior auditor or increasing the monitoring frequency from the start.", icon="❓")
        with st.form("risk_predictor_form"):
            col1, col2, col3 = st.columns(3)
            with col1: p_type = st.selectbox("Trial Type", portfolio_df['Trial_Type'].unique(), key='p_type')
            with col2: p_phase = st.selectbox("Trial Phase", portfolio_df['Phase'].unique(), key='p_phase')
            with col3: p_pi_exp = st.selectbox("PI Experience", portfolio_df['PI_Experience_Level'].unique(), key='p_pi')
            if st.form_submit_button("🔬 Forecast Risk Profile"):
                input_df = pd.DataFrame({'Trial_Type': [p_type], 'Phase': [p_phase], 'PI_Experience_Level': [p_pi_exp], 'Is_First_In_Human': [False], 'Num_Sites': [1]})
                input_encoded = pd.DataFrame(encoder.transform(input_df[['Trial_Type', 'Phase', 'PI_Experience_Level']]), columns=encoder.get_feature_names_out(['Trial_Type', 'Phase', 'PI_Experience_Level']))
                input_final = pd.concat([input_df.drop(columns=['Trial_Type', 'Phase', 'PI_Experience_Level']).reset_index(drop=True), input_encoded], axis=1).reindex(columns=model_features, fill_value=0)
                prediction_proba = risk_model.predict_proba(input_final)[0][1]
                st.success(f"Predicted Risk of a Major Finding: **{prediction_proba:.1%}**")

    with st.container(border=True):
        st.markdown("##### Key Risk Factor Analysis")
        st.info("💡 **Expert Tip:** This chart shows which trial characteristics have the biggest impact on the likelihood of a major finding. Use this to justify resource allocation and guide preventative training efforts.", icon="❓")
        importance_df['Color'] = np.where(importance_df['Coefficient'] < 0, '#005A9C', '#E63946')
        fig = go.Figure(go.Bar(x=importance_df.index, y=importance_df['Coefficient'], marker_color=importance_df['Color']))
        fig.update_layout(title_text="<b>Impact of Trial Factors on Risk of Major Findings</b>", xaxis_title="Trial Characteristic", yaxis_title="Model Coefficient (Impact on Risk)")
        st.plotly_chart(fig, use_container_width=True)

    with st.container(border=True):
        st.markdown("##### Portfolio Risk Segmentation")
        st.info("💡 **Expert Tip:** This sunburst chart shows how the overall portfolio risk is distributed across different trial types and disease teams. A large 'High Risk' segment in a specific area may indicate a need for a targeted review.", icon="❓")
        X_all_cat = portfolio_df[['Trial_Type', 'Phase', 'PI_Experience_Level']]
        X_all_num = portfolio_df[['Is_First_In_Human', 'Num_Sites']]
        X_all_encoded = pd.DataFrame(encoder.transform(X_all_cat), columns=encoder.get_feature_names_out(X_all_cat.columns))
        X_all_final = pd.concat([X_all_num.reset_index(drop=True), X_all_encoded], axis=1).reindex(columns=model_features, fill_value=0)
        portfolio_df['Predicted_Risk'] = risk_model.predict_proba(X_all_final)[:, 1]
        portfolio_df['Risk_Tier'] = pd.cut(portfolio_df['Predicted_Risk'], bins=[0, 0.2, 0.4, 0.6, 1.0], labels=['Low', 'Medium', 'High', 'Critical'], right=False)
        path_df = portfolio_df.groupby(['Risk_Tier', 'Disease_Team', 'Trial_Type']).size().reset_index(name='Count')
        fig = px.sunburst(path_df, path=['Risk_Tier', 'Disease_Team', 'Trial_Type'], values='Count', title="<b>Portfolio Risk Distribution</b>")
        fig.update_layout(margin=dict(t=40, l=0, r=0, b=0))
        st.plotly_chart(fig, use_container_width=True)

def render_systemic_risk(findings_df, portfolio_df):
    st.subheader("Systemic Process & Risk Analysis", divider="blue")
    st.markdown("_This section moves beyond individual data points to identify systemic trends, process vulnerabilities, and non-random patterns across the clinical trial portfolio._")
    col1, col2 = st.columns(2)
    with col1:
        with st.container(border=True):
            st.markdown("##### Statistical Process Control (SPC) Analysis")
            st.info("💡 **Expert Tip:** Any red 'X' points outside the shaded blue area represent 'special cause' variations that require immediate Root Cause Analysis.", icon="❓")
            category_to_monitor = st.selectbox("Select Finding Category to Analyze for Trends:", options=findings_df['Category'].unique())
            st.plotly_chart(plot_spc_chart_sme(findings_df, 'Finding_Date', 'Category', category_to_monitor, f"SPC for '{category_to_monitor}'"), use_container_width=True)
    with col2:
        with st.container(border=True):
            st.markdown("##### Finding Concentration Analysis")
            st.info("💡 **Expert Tip:** Dark blue 'hotspots' suggest a targeted training intervention for that specific team is more effective than a CTO-wide initiative.", icon="❓")
            st.plotly_chart(px.imshow(pd.crosstab(findings_df['Disease_Team'], findings_df['Category']), text_auto=True, aspect="auto", title="<b>Finding Concentration by Disease Team & Category</b>", color_continuous_scale='Blues'), use_container_width=True)
    
    with st.container(border=True):
        st.markdown("##### CAPA Effectiveness & Finding Recurrence Analysis")
        st.info("💡 **Expert Tip:** This chart is critical for a mature QMS. A high recurrence rate in a specific category indicates that our Corrective and Preventive Actions are not addressing the root cause and need to be re-evaluated.", icon="❓")
        capa_effectiveness_df = analyze_capa_effectiveness(findings_df)
        if not capa_effectiveness_df.empty:
            fig = px.bar(capa_effectiveness_df, x='Recurrence_Rate', y='Category', orientation='h', title='<b>Finding Recurrence Rate by Category (Post-CAPA)</b>', labels={'Recurrence_Rate': 'Recurrence Rate (%)', 'Category': 'Finding Category'}, text=capa_effectiveness_df['Recurrence_Rate'].apply(lambda x: f'{x:.0%}'))
            fig.update_traces(marker_color='#E63946', textposition='outside')
            fig.update_layout(plot_bgcolor='white', xaxis_tickformat='.0%', yaxis={'categoryorder':'total ascending'})
            st.plotly_chart(fig, use_container_width=True)
        else: st.write("No closed CAPAs available to analyze for effectiveness.")

    with st.container(border=True):
        st.markdown("##### Interactive Regulatory Inspection Simulation")
        st.info("💡 **Expert Tip:** Use this tool to pressure-test readiness and train staff. Can the requested document be produced instantly? This simulates the high-stakes reality of an FDA or EMA audit.", icon="❓")
        @st.cache_data
        def create_mock_etmf(_portfolio_df):
            mock_etmf = {}
            for _, row in _portfolio_df.iterrows():
                trial_id, num_subjects = row['Trial_ID'], row['Subjects_Enrolled']
                mock_etmf[trial_id] = {"Protocol Signature Page": f"DocRef_PSP_{trial_id}.pdf", "IRB Approval Letter": f"DocRef_IRB_Approval_{trial_id}.pdf", "FDA Form 1572": f"DocRef_1572_{trial_id}.pdf" if "IIT" in trial_id or "IND" in trial_id else "N/A", "Informed Consent Forms": {f"Subject-{i:03d}": f"ICF_{trial_id}_Subj_{i:03d}.pdf" for i in range(1, num_subjects + 1)}, "Serious Adverse Event Reports": {f"SAE-{i:03d}": f"SAE_{trial_id}_{i:03d}.pdf" for i in range(1, np.random.randint(2, 6))}}
            return mock_etmf
        mock_etmf_db = create_mock_etmf(portfolio_df)
        sim_col1, sim_col2 = st.columns([1, 2])
        with sim_col1:
            st.write("**Inspection Scenario:**")
            trial_to_inspect = st.selectbox("Select a Trial to Inspect:", options=portfolio_df['Trial_ID'], key="inspect_trial")
            subject_list = list(mock_etmf_db[trial_to_inspect]["Informed Consent Forms"].keys())
            if subject_list: subject_to_inspect = st.selectbox("Select a Subject:", options=subject_list, key="inspect_subject")
            else: subject_to_inspect = None
        if subject_to_inspect and st.button("🔬 Pull Subject's Consent Form"):
            with sim_col2:
                st.info(f"Request: 'Please provide the signed consent form for {subject_to_inspect} on trial {trial_to_inspect}.'")
                with st.spinner("Searching eTMF..."): time.sleep(1.5)
                st.success("**Document Found!**")
                st.code(f"File Path: /eTMF/Trials/{trial_to_inspect}/Subject_Files/{subject_to_inspect}/ICF_{trial_to_inspect}_Subj_{subject_to_inspect.split('-')[1]}.pdf\nAccessed: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", language="bash")
        if st.button("📄 Pull Trial's 1572 Form"):
            with sim_col2:
                st.info(f"Request: 'Please provide the current FDA Form 1572 for trial {trial_to_inspect}.'")
                with st.spinner("Searching eTMF..."): time.sleep(1)
                doc_ref = mock_etmf_db[trial_to_inspect].get("FDA Form 1572")
                if doc_ref != "N/A": 
                    st.success("**Document Found!**")
                    st.code(f"File Path: /eTMF/Trials/{trial_to_inspect}/Regulatory/{doc_ref}\nAccessed: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", language="bash")
                else: 
                    st.error("**Document Not Applicable!** This is not an IND-holding trial.")

def render_pi_performance(portfolio_df, findings_df):
    st.subheader("Principal Investigator (PI) Performance Oversight", divider="blue")
    st.markdown("_This section provides tools for both team-level oversight and individual PI coaching, enabling targeted support and recognition._")

    all_teams = sorted(portfolio_df['Disease_Team'].unique())
    selected_team = st.selectbox("Select a Disease Team to Analyze:", options=all_teams, key="team_select")
    
    team_trials = portfolio_df[portfolio_df['Disease_Team'] == selected_team]
    team_findings = findings_df[findings_df['Disease_Team'] == selected_team]
    
    # --- 1. Team-Level Quality Scorecard ---
    with st.container(border=True):
        st.markdown(f"##### 🏆 Disease Team Quality Scorecard: **{selected_team}**")
        st.info("💡 **Expert Tip:** Use this scorecard for Disease Team meetings to identify top performers for recognition and investigators who may need additional support or resources.", icon="❓")
        
        if not team_trials.empty:
            # Aggregate metrics by PI
            pi_metrics = team_trials.groupby('PI_Name').agg(
                Trials_Open=('Trial_ID', 'count'),
                Total_Enrolled=('Subjects_Enrolled', 'sum'),
                Avg_Accrual=('Avg_Accrual_Per_Month', 'mean')
            ).reset_index()

            pi_findings_summary = team_findings.groupby('PI_Name').agg(
                Critical_Findings=('Risk_Level', lambda x: (x == 'Critical').sum()),
                Major_Findings=('Risk_Level', lambda x: (x == 'Major').sum()),
                Overdue_CAPAs=('CAPA_Status', lambda x: (x == 'Overdue').sum()),
                Avg_Closure_Days=('Days_to_Close', 'mean')
            ).reset_index()
            
            # Merge into a final scorecard
            scorecard_df = pd.merge(pi_metrics, pi_findings_summary, on='PI_Name', how='left').fillna(0)
            scorecard_df = scorecard_df.sort_values(by=['Critical_Findings', 'Major_Findings', 'Overdue_CAPAs'], ascending=[False, False, False])
            
            st.dataframe(
                scorecard_df.style.format({
                    'Avg_Accrual': '{:.1f}',
                    'Avg_Closure_Days': '{:.1f}'
                }).background_gradient(
                    cmap='Reds', subset=['Critical_Findings', 'Major_Findings', 'Overdue_CAPAs']
                ).background_gradient(
                    cmap='Greens', subset=['Avg_Accrual']
                ).bar(subset=['Total_Enrolled'], color='#5cadff'),
                use_container_width=True,
                hide_index=True
            )
        else:
            st.warning("No trial data available for this team.")

    # --- 2. Individual PI Drill-Down ---
    with st.container(border=True):
        st.markdown("##### 🧑‍🔬 Individual Investigator Drill-Down")
        available_pis = sorted(team_trials['PI_Name'].unique())
        if available_pis:
            selected_pi = st.selectbox("Select a PI for Detailed Analysis:", options=available_pis, key="pi_select")
            
            pi_findings = findings_df[findings_df['PI_Name'] == selected_pi]
            
            # The individual snapshot and barchart from the original function
            m_col1, m_col2, m_col3 = st.columns(3)
            pi_major_findings = pi_findings[pi_findings['Risk_Level'].isin(['Major', 'Critical'])].shape[0]
            team_major_findings_grouped = team_findings[team_findings['Risk_Level'].isin(['Major', 'Critical'])].groupby('PI_Name').size()
            team_avg_major_findings = team_major_findings_grouped.mean() if not team_major_findings_grouped.empty else 0
            m_col1.metric("Major/Critical Findings", f"{pi_major_findings}", f"Team Avg: {team_avg_major_findings:.1f}", delta_color="inverse")
            
            pi_avg_closure = pi_findings['Days_to_Close'].mean()
            team_avg_closure = team_findings['Days_to_Close'].mean()
            m_col2.metric("Avg. CAPA Closure (Days)", f"{pi_avg_closure:.1f}" if not np.isnan(pi_avg_closure) else "N/A", f"Team Avg: {team_avg_closure:.1f}", delta_color="inverse")
            
            pi_overdue = pi_findings[pi_findings['CAPA_Status'] == 'Overdue'].shape[0]
            team_overdue_grouped = team_findings.groupby('PI_Name')['CAPA_Status'].apply(lambda x: (x == 'Overdue').sum())
            team_avg_overdue = team_overdue_grouped.mean() if not team_overdue_grouped.empty else 0
            m_col3.metric("Overdue CAPAs", f"{pi_overdue}", f"Team Avg: {team_avg_overdue:.1f}", delta_color="inverse")

            st.plotly_chart(plot_pi_findings_barchart_sme(pi_findings, selected_pi), use_container_width=True)
        else:
            st.info("No PIs found for the selected team.")

def render_organizational_capability(team_df, initiatives_df, audits_df, findings_df):
    st.subheader("Organizational Capability & Strategic Oversight", divider="blue")
    st.markdown("_This section assesses the capacity of the QA team and tracks progress and financial health of key strategic objectives._")

    main_tabs = st.tabs(["📊 Team Performance & Workload", "🚀 Strategic Initiatives"])
    
    with main_tabs[0]:
        st.markdown("##### Auditor Performance, Skills, and Process Efficiency")
        
        auditor_options = ['All Auditors'] + sorted(team_df['Auditor'].unique())
        selected_auditors = st.multiselect(
            "Select Auditors to Analyze:",
            options=auditor_options,
            default=['All Auditors']
        )

        if 'All Auditors' in selected_auditors or not selected_auditors:
            filtered_team_df = team_df.copy()
            filtered_audit_yield_df = audit_yield_df.copy()
        else:
            filtered_team_df = team_df[team_df['Auditor'].isin(selected_auditors)]
            filtered_audit_yield_df = audit_yield_df[audit_yield_df['Auditor'].isin(selected_auditors)]

        plot_tabs = st.tabs(["🎯 Skill & Performance Quadrant", "📈 Efficiency & Workload Trends", "🔍 Audit Yield Analysis", "✅ Compliance & Readiness"])

        with plot_tabs[0]:
            st.markdown("###### Skill & Performance Quadrant Analysis")
            st.info("💡 **Expert Tip:** This chart combines skill with performance. Use it to identify auditors who are both highly skilled and efficient (top-right quadrant) as potential mentors, or those who may need coaching (bottom-left).", icon="❓")
            
            skill_cols = [col for col in filtered_team_df.columns if '_Skill' in col]
            filtered_team_df['Avg_Skill_Score'] = filtered_team_df[skill_cols].mean(axis=1)
            
            if 'Strain' not in filtered_team_df.columns:
                filtered_team_df['Skill_Factor'] = filtered_team_df['IIT_Oversight_Skill'] + filtered_team_df['FDA_Inspection_Mgmt_Skill']
                filtered_team_df['Strain'] = (filtered_team_df['Audits_Conducted_YTD'] * filtered_team_df['Avg_Report_Turnaround_Days']) / (filtered_team_df['Skill_Factor'] + 1)

            fig = px.scatter(
                filtered_team_df,
                x='Avg_Skill_Score',
                y='Avg_Report_Turnaround_Days',
                size='Audits_Conducted_YTD',
                color='Strain',
                text='Auditor',
                title='<b>Auditor Skill vs. Efficiency Quadrant</b>',
                labels={
                    'Avg_Skill_Score': 'Average Skill Score (All Categories)',
                    'Avg_Report_Turnaround_Days': 'Efficiency (Lower is Better)',
                    'Audits_Conducted_YTD': 'Workload (Audits YTD)',
                    'Strain': 'Strain Index'
                },
                color_continuous_scale=px.colors.sequential.RdYlGn_r # Reversed: Red is high strain
            )
            
            fig.update_yaxes(autorange="reversed")
            
            mean_x = filtered_team_df['Avg_Skill_Score'].mean()
            mean_y = filtered_team_df['Avg_Report_Turnaround_Days'].mean()
            
            fig.add_hline(y=mean_y, line_dash="dot", line_color="grey", annotation_text="Team Avg. Efficiency")
            fig.add_vline(x=mean_x, line_dash="dot", line_color="grey", annotation_text="Team Avg. Skill")
            
            fig.update_traces(textposition='top center')
            fig.update_layout(plot_bgcolor='white', height=500)
            st.plotly_chart(fig, use_container_width=True)

        with plot_tabs[1]:
            st.markdown("###### Historical Workload vs. Efficiency Trends")
            st.info("💡 **Expert Tip:** Watch for divergence. If the blue workload line rises while the red efficiency line also rises (gets worse), it's a strong indicator of impending team burnout or process bottlenecks.", icon="❓")
            
            trend_df = filtered_audit_yield_df.copy()
            trend_df['Quarter'] = pd.to_datetime(trend_df['Audit_Date']).dt.to_period('Q').astype(str)
            
            quarterly_summary = trend_df.groupby('Quarter').agg(
                Avg_Turnaround=('Turnaround_Time', 'mean'),
                Total_Audits=('Audit_ID', 'count')
            ).reset_index()

            fig = go.Figure()
            fig.add_trace(go.Bar(
                x=quarterly_summary['Quarter'],
                y=quarterly_summary['Total_Audits'],
                name='Audits Conducted',
                marker_color='#A8DADC'
            ))
            fig.add_trace(go.Scatter(
                x=quarterly_summary['Quarter'],
                y=quarterly_summary['Avg_Turnaround'],
                name='Avg. Turnaround (Days)',
                yaxis='y2',
                mode='lines+markers',
                line=dict(color='#E63946', width=3)
            ))
            
            fig.update_layout(
                title_text="<b>Quarterly Workload vs. Report Turnaround Efficiency</b>",
                yaxis=dict(title='Total Audits Conducted'),
                yaxis2=dict(title='Avg. Turnaround (Days)', overlaying='y', side='right', showgrid=False, autorange="reversed"),
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                plot_bgcolor='white'
            )
            st.plotly_chart(fig, use_container_width=True)

        with plot_tabs[2]:
            st.markdown("###### Comparative Audit Yield Analysis")
            st.info("💡 **Expert Tip:** This chart shows the 'quality' of audits. Are some auditors finding significantly more critical issues than their peers? This can guide training, calibration sessions, or future audit assignments.", icon="❓")
            
            avg_yield_by_auditor = filtered_audit_yield_df.groupby('Auditor')[['Critical', 'Major', 'Minor']].mean()
            team_avg = avg_yield_by_auditor.mean().rename('Team Average')
            
            plot_data = pd.concat([avg_yield_by_auditor, pd.DataFrame(team_avg).T])

            fig = go.Figure()
            colors = {'Critical': '#E63946', 'Major': '#FFC72C', 'Minor': '#457B9D'}
            for risk_level in ['Critical', 'Major', 'Minor']:
                fig.add_trace(go.Bar(
                    x=plot_data.index,
                    y=plot_data[risk_level],
                    name=risk_level,
                    marker_color=colors[risk_level]
                ))

            fig.update_layout(
                barmode='stack',
                title_text="<b>Average Findings per Audit by Auditor vs. Team Average</b>",
                yaxis_title="Avg. Findings per Audit",
                xaxis_title=None,
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                plot_bgcolor='white'
            )
            st.plotly_chart(fig, use_container_width=True)
            
        with plot_tabs[3]:
            st.markdown("###### QA Team Compliance & Readiness")
            st.info("💡 **Expert Tip:** This provides an immediate view of departmental compliance risk. An expired certification is a direct liability during an inspection of the QA unit itself.", icon="❓")
            
            col1, col2 = st.columns(2)
            with col1:
                compliance_df = filtered_team_df['GCP_Certification_Status'].value_counts().reset_index()
                compliance_df.columns = ['Status', 'Count']
                
                fig = go.Figure(go.Indicator(
                    mode="gauge+number",
                    value=len(filtered_team_df[filtered_team_df['GCP_Certification_Status'] == 'Current']),
                    number={'suffix': f" / {len(filtered_team_df)}"},
                    domain={'x': [0, 1], 'y': [0, 1]},
                    title={'text': "Auditors with Current GCP Certification"},
                    gauge={'axis': {'range': [None, len(filtered_team_df)]},
                           'bar': {'color': "#457B9D"},
                           'steps': [
                               {'range': [0, len(filtered_team_df) * 0.8], 'color': "#E63946"},
                               {'range': [len(filtered_team_df) * 0.8, len(filtered_team_df) * 0.95], 'color': "#FFC72C"},
                           ]}
                ))
                fig.update_layout(height=250, margin=dict(t=40, b=40))
                st.plotly_chart(fig, use_container_width=True)

            with col2:
                expiring_soon = filtered_team_df[filtered_team_df['GCP_Certification_Status'] == 'Expires <90d']
                if not expiring_soon.empty:
                    st.warning("Certifications Expiring Soon:")
                    st.dataframe(expiring_soon[['Auditor']], use_container_width=True, hide_index=True)
                else:
                    st.success("All team certifications are current for the next 90 days.")

    with main_tabs[1]:
        st.markdown("##### Strategic Initiatives & Financial Oversight")
        st.info("💡 **Expert Tip:** A CPI or SPI value < 1.0 indicates a project is over budget or behind schedule, respectively. This allows for proactive intervention before projects go significantly off-track.", icon="❓")
        today = pd.to_datetime(datetime.datetime.now())
        initiatives_df['Start_Date'] = pd.to_datetime(initiatives_df['Start_Date'])
        initiatives_df['End_Date'] = pd.to_datetime(initiatives_df['End_Date'])
        initiatives_df['Days_Elapsed'] = (today - initiatives_df['Start_Date']).dt.days
        initiatives_df['Total_Days_Planned'] = (initiatives_df['End_Date'] - initiatives_df['Start_Date']).dt.days
        initiatives_df['Daily_Burn_Rate'] = initiatives_df.apply(lambda row: row['Spent_USD'] / row['Days_Elapsed'] if row['Days_Elapsed'] > 0 else 0, axis=1)
        initiatives_df['Projected_Total_Cost'] = initiatives_df['Daily_Burn_Rate'] * initiatives_df['Total_Days_Planned']
        initiatives_df['Projected_Over_Under'] = initiatives_df['Budget_USD'] - initiatives_df['Projected_Total_Cost']
        initiatives_df['CPI'] = initiatives_df.apply(lambda row: (row['Budget_USD'] * (row['Percent_Complete']/100)) / row['Spent_USD'] if row['Spent_USD'] > 0 else 0, axis=1)
        initiatives_df['SPI'] = initiatives_df.apply(lambda row: (row['Percent_Complete']/100) / (row['Days_Elapsed']/row['Total_Days_Planned']) if row['Days_Elapsed'] > 0 and row['Total_Days_Planned'] > 0 else 0, axis=1)
        
        def format_financials(df):
            return df.style.format({
                'Budget_USD': "${:,.0f}", 'Spent_USD': "${:,.0f}", 
                'Projected_Total_Cost': "${:,.0f}", 'Projected_Over_Under': "${:,.0f}", 
                'Daily_Burn_Rate': "${:,.2f}", 'CPI': "{:.2f}", 'SPI': "{:.2f}"
            }).background_gradient(cmap='RdYlGn', subset=['Projected_Over_Under']
            ).background_gradient(cmap='RdYlGn', vmin=0.8, vmax=1.2, subset=['CPI', 'SPI']
            ).bar(subset=['Percent_Complete'], color='#5cadff', vmin=0, vmax=100)
            
        st.dataframe(format_financials(initiatives_df[['Initiative', 'Lead', 'Status', 'Percent_Complete', 'Budget_USD', 'Spent_USD', 'Projected_Total_Cost', 'Projected_Over_Under', 'CPI', 'SPI']]), use_container_width=True)
        st.caption("CPI (Cost Performance Index) & SPI (Schedule Performance Index): > 1.0 is favorable (green), < 1.0 is unfavorable (red).")
        
        st.markdown("---")
        st.markdown("###### 🧪 Initiative Planner & Budget Simulator")
        st.info("💡 **Expert Tip:** Model a new QA initiative to understand its potential financial impact before committing resources. This allows for data-driven budget requests.", icon="❓")

        with st.form("budget_simulator_form"):
            s_col1, s_col2, s_col3 = st.columns(3)
            with s_col1:
                new_initiative_name = st.text_input("Initiative Name", "New eRegulatory System Rollout")
                new_budget = st.number_input("Proposed Budget (USD)", min_value=5000, max_value=250000, value=50000, step=1000)
            with s_col2:
                new_start_date = st.date_input("Proposed Start Date", datetime.date.today())
                new_end_date = st.date_input("Proposed End Date", datetime.date.today() + datetime.timedelta(days=365))
            with s_col3:
                num_personnel = st.slider("Dedicated Personnel (FTE)", min_value=0.25, max_value=5.0, value=1.5, step=0.25)
                avg_salary = 110000 # Assume an average burdened salary for planning
                
            submitted = st.form_submit_button("Simulate Financials")

        if submitted:
            if new_end_date <= new_start_date:
                st.error("End Date must be after Start Date.")
            else:
                total_days = (new_end_date - new_start_date).days
                personnel_cost = (total_days / 365) * num_personnel * avg_salary
                discretionary_budget = new_budget - personnel_cost
                
                st.success(f"**Simulation Results for '{new_initiative_name}'**")
                sim_col1, sim_col2, sim_col3 = st.columns(3)
                sim_col1.metric("Total Project Duration", f"{total_days} days")
                sim_col2.metric("Projected Personnel Cost", f"${personnel_cost:,.0f}")
                sim_col3.metric("Remaining Discretionary Budget", f"${discretionary_budget:,.0f}", help="Funds available for software, training, etc., after accounting for salaries.")

def render_quality_impact(portfolio_df, findings_df):
    st.subheader("Quality Impact on Trial Performance", divider="blue")
    st.markdown("_This analysis correlates quality metrics with key trial operational metrics to demonstrate the tangible impact of quality on research success._")
    
    major_findings_count = findings_df[findings_df['Risk_Level'].isin(['Major', 'Critical'])].groupby('Trial_ID').size().reset_index(name='Major_Finding_Count')
    plot_df = pd.merge(portfolio_df, major_findings_count, on='Trial_ID', how='left').fillna(0)
    
    tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(["📈 Accrual vs. Findings", "📉 Screen Failures", "💻 Data Quality", "⏳ Startup Delays", "🧑‍🏫 PI Experience", " holistic view The Big Picture"])
    with tab1:
        st.markdown("##### Correlation: Major Findings vs. Patient Accrual")
        st.info("💡 **Expert Tip:** Use this chart to demonstrate to Disease Team Leaders and PIs how quality issues can directly impact trial performance and enrollment goals. Trials with fewer major findings (left side) tend to have higher monthly accrual.", icon="❓")
        fig = px.scatter(plot_df, x='Major_Finding_Count', y='Avg_Accrual_Per_Month', size='Subjects_Enrolled', color='Disease_Team', hover_name='Trial_ID', title='<b>Higher Quality Correlates with Faster Patient Accrual</b>', labels={'Major_Finding_Count': 'Number of Critical/Major Findings', 'Avg_Accrual_Per_Month': 'Average Accrual per Month'}, size_max=40)
        fig.update_layout(plot_bgcolor='white')
        st.plotly_chart(fig, use_container_width=True)
    with tab2:
        st.markdown("##### Impact of Consent & Eligibility Findings on Screen Failures")
        st.info("💡 **Expert Tip:** This plot provides direct evidence that poor upfront quality control in consent and eligibility leads to wasted resources (higher screen fail rates). Use this to justify investment in better pre-screening and consent monitoring processes.", icon="❓")
        eligibility_findings = findings_df[findings_df['Category'].isin(['Informed Consent Process', 'Protocol Adherence'])].groupby('Trial_ID').size().reset_index(name='Eligibility_Finding_Count')
        plot_df_sfr = pd.merge(portfolio_df, eligibility_findings, on='Trial_ID', how='left').fillna(0)
        fig = px.scatter(plot_df_sfr, x='Eligibility_Finding_Count', y='Screen_Fail_Rate', color='Phase', size='Subjects_Enrolled', hover_name='Trial_ID', title='<b>More Eligibility-Related Findings Correlate with Higher Screen Fail Rates</b>', labels={'Eligibility_Finding_Count': 'Consent & Eligibility Findings', 'Screen_Fail_Rate': 'Screen Failure Rate'})
        fig.update_layout(plot_bgcolor='white', yaxis_tickformat=".0%")
        st.plotly_chart(fig, use_container_width=True)
    with tab3:
        st.markdown("##### Impact of Source Data Findings on Data Query Rate")
        st.info("💡 **Expert Tip:** This chart links QA oversight directly to data management efficiency. Show this to data management leaders to build a partnership around improving source data quality, which reduces their workload and accelerates time to database lock.", icon="❓")
        sdv_findings = findings_df[findings_df['Category'] == 'Source Data Verification'].groupby('Trial_ID').size().reset_index(name='SDV_Finding_Count')
        plot_df_query = pd.merge(portfolio_df, sdv_findings, on='Trial_ID', how='left').fillna(0)
        fig = px.scatter(plot_df_query, x='SDV_Finding_Count', y='Data_Query_Rate', color='Trial_Type', size='Subjects_Enrolled', hover_name='Trial_ID', title='<b>More Source Data Findings Correlate with Higher Data Query Rates</b>', labels={'SDV_Finding_Count': 'Number of Source Data Verification Findings', 'Data_Query_Rate': 'Data Queries per Subject'})
        fig.update_layout(plot_bgcolor='white', yaxis_tickformat=".1%")
        st.plotly_chart(fig, use_container_width=True)
    with tab4:
        st.markdown("##### Impact of Regulatory Findings on Trial Startup Times")
        st.info("💡 **Expert Tip:** This chart demonstrates that early, robust quality checks in the startup phase can reduce activation delays. Use this to advocate for QA involvement *before* a trial is activated, not just after.", icon="❓")
        startup_findings = findings_df[findings_df['Category'] == 'Regulatory Binder Mgmt'].groupby('Trial_ID').size().reset_index(name='Startup_Finding_Count')
        plot_df_startup = pd.merge(portfolio_df, startup_findings, on='Trial_ID', how='left').fillna(0)
        fig = px.scatter(plot_df_startup, x='Startup_Finding_Count', y='Days_to_Activation', color='PI_Experience_Level', hover_name='Trial_ID', title='<b>More Regulatory Findings Correlate with Longer Startup Times</b>', labels={'Startup_Finding_Count': 'Number of Startup-Phase Findings', 'Days_to_Activation': 'Days from IRB Submission to Activation'})
        fig.update_layout(plot_bgcolor='white')
        st.plotly_chart(fig, use_container_width=True)
    with tab5:
        st.markdown("##### PI Experience vs. Quality & Performance")
        st.info("💡 **Expert Tip:** This provides objective evidence for the 'New Investigator Program.' It shows that while new PIs have more quality findings, they can still achieve high accrual. This justifies providing them with enhanced, targeted QA support to mitigate risk without stifling their research.", icon="❓")
        fig = px.scatter(plot_df, x='Avg_Accrual_Per_Month', y='Major_Finding_Count', color='PI_Experience_Level', size='Subjects_Enrolled', title='<b>Relationship Between PI Experience, Quality, and Accrual</b>', labels={'Avg_Accrual_Per_Month': 'Average Accrual per Month', 'Major_Finding_Count': 'Number of Major/Critical Findings'}, color_discrete_map={'New': '#E63946', 'Intermediate': '#FFC72C', 'Expert': '#005A9C'})
        fig.update_layout(plot_bgcolor='white')
        st.plotly_chart(fig, use_container_width=True)
    with tab6:
        st.markdown("##### The Big Picture: Cost of Poor Quality")
        st.info("💡 **Expert Tip:** This is your executive summary. It aggregates the data to make a powerful business case for investing in quality. It clearly shows that higher-risk trials (as defined by our predictive model) are slower, less efficient, and costlier across the board.", icon="❓")
        risk_model, encoder, model_features, _ = get_risk_model_with_importance(portfolio_df)
        X_all_cat = plot_df[['Trial_Type', 'Phase', 'PI_Experience_Level']]
        X_all_num = plot_df[['Is_First_In_Human', 'Num_Sites']]
        X_all_encoded = pd.DataFrame(encoder.transform(X_all_cat), columns=encoder.get_feature_names_out(X_all_cat.columns))
        X_all_final = pd.concat([X_all_num.reset_index(drop=True), X_all_encoded], axis=1).reindex(columns=model_features, fill_value=0)
        plot_df['Predicted_Risk'] = risk_model.predict_proba(X_all_final)[:, 1]
        plot_df['Risk_Tier'] = pd.cut(plot_df['Predicted_Risk'], bins=[0, 0.4, 0.6, 1.0], labels=['Low Risk', 'Medium Risk', 'High Risk'], right=False)
        impact_summary = plot_df.groupby('Risk_Tier', observed=True)[['Avg_Accrual_Per_Month', 'Screen_Fail_Rate', 'Data_Query_Rate']].mean(numeric_only=True).reset_index()
        fig = go.Figure(data=[
            go.Bar(name='Avg Accrual/Month', x=impact_summary['Risk_Tier'], y=impact_summary['Avg_Accrual_Per_Month'], yaxis='y', offsetgroup=1),
            go.Bar(name='Screen Fail Rate', x=impact_summary['Risk_Tier'], y=impact_summary['Screen_Fail_Rate'], yaxis='y2', offsetgroup=2),
            go.Bar(name='Data Query Rate', x=impact_summary['Risk_Tier'], y=impact_summary['Data_Query_Rate'], yaxis='y2', offsetgroup=2, base=impact_summary['Screen_Fail_Rate'])
        ],
        layout={'yaxis': {'title': 'Avg Accrual per Month'}, 'yaxis2': {'title': 'Rate (%)', 'overlaying': 'y', 'side': 'right', 'tickformat': '.0%'}, 'title': '<b>The Operational Cost of Poor Quality by Predicted Risk Tier</b>'})
        fig.update_layout(barmode='group')
        st.plotly_chart(fig, use_container_width=True)

# ======================================================================================
# SECTION 5: MAIN APP ORCHESTRATION
# ======================================================================================
def main():
    with st.sidebar:
        st.markdown("## Moores Cancer Center")
        st.markdown("### Clinical Trials Office")
        st.markdown("---")
        selected = option_menu(
            menu_title="QA Command Center",
            options=["Home", "Predictive Analytics", "Systemic Risk", "PI Performance", "Team & Strategy", "Quality Impact"],
            icons=["house-door-fill", "graph-up-arrow", "shield-shaded", "person-badge", "people-fill", "broadcast-pin"],
            menu_icon="kanban-fill", default_index=0
        )
        st.markdown("---")
        st.info("This dashboard is a prototype demonstrating a proactive, data-driven approach to Clinical QA management.")
        st.header("Generate Executive Report")
        st.info("Download a PowerPoint summary of the current QA program status for leadership review.")
        
        portfolio_df_sidebar, findings_df_sidebar, team_df_sidebar, _, _ = generate_master_data()
       
        risk_weights = {'Critical': 10, 'Major': 5, 'Minor': 1}
        findings_df_sidebar['Risk_Score'] = findings_df_sidebar['Risk_Level'].map(risk_weights)
        open_findings_sidebar = findings_df_sidebar[~findings_df_sidebar['CAPA_Status'].isin(['Closed-Effective'])].copy()
        
        total_risk_score_sidebar = int(open_findings_sidebar['Risk_Score'].sum())
        overdue_major_capas_sidebar = findings_df_sidebar[(findings_df_sidebar['CAPA_Status'] == 'Overdue') & (findings_df_sidebar['Risk_Level'] != 'Minor')].shape[0]
        readiness_score_sidebar = max(0, 100 - (overdue_major_capas_sidebar * 10) - (open_findings_sidebar[open_findings_sidebar['Risk_Level'] == 'Critical'].shape[0] * 5))
        team_df_sidebar['Skill_Factor'] = team_df_sidebar['IIT_Oversight_Skill'] + team_df_sidebar['FDA_Inspection_Mgmt_Skill']
        team_df_sidebar['Strain'] = (team_df_sidebar['Audits_Conducted_YTD'] * team_df_sidebar['Avg_Report_Turnaround_Days']) / (team_df_sidebar['Skill_Factor'] + 1)
        avg_strain_sidebar = team_df_sidebar['Strain'].mean()
        
        kpi_data_for_report = [("Portfolio Risk Score", total_risk_score_sidebar, f"{open_findings_sidebar[open_findings_sidebar['Risk_Level'] == 'Critical'].shape[0]} Open Criticals"), ("Inspection Readiness", f"{readiness_score_sidebar}%", f"{overdue_major_capas_sidebar} Overdue Major CAPAs"), ("Resource Strain Index", f"{avg_strain_sidebar:.2f}", "Target < 2.5")]
        findings_for_report = findings_df_sidebar[(findings_df_sidebar['Risk_Level'].isin(['Major', 'Critical'])) & (findings_df_sidebar['CAPA_Status'] != 'Closed-Effective')][['Trial_ID', 'Category', 'Risk_Level', 'CAPA_Status']].head(10)
        default_spc_fig = plot_spc_chart_sme(findings_df_sidebar, 'Finding_Date', 'Category', 'Informed Consent Process', "SPC Chart: Informed Consent Findings")
        
        ppt_buffer = generate_ppt_report(kpi_data_for_report, default_spc_fig, findings_for_report)
        st.download_button(label="📥 Download PowerPoint Report", data=ppt_buffer, file_name=f"MCC_CTO_QA_Summary_{datetime.date.today()}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        st.markdown("---")
        st.markdown("### Key Concepts & Regulations")
        st.markdown("- **RBQM:** Risk-Based Quality Management\n- **SPC:** Statistical Process Control\n- **CPI/SPI:** Cost/Schedule Performance Index\n- **GCP:** Good Clinical Practice\n- **21 CFR Part 50 & 312:** Key FDA Regulations")

    st.title("🔬 Scientific QA Command Center")
    st.markdown("An advanced analytics dashboard for the Assistant Director of Quality Assurance.")

    portfolio_df, findings_df, team_df, initiatives_df, audits_df = generate_master_data()
    
    findings_df['Closure_Date'] = findings_df.apply(lambda row: row['Finding_Date'] + pd.to_timedelta(np.random.randint(5, 60), unit='d') if row['CAPA_Status'] == 'Closed-Effective' else pd.NaT, axis=1)
    findings_df['Days_to_Close'] = (findings_df['Closure_Date'] - findings_df['Finding_Date']).dt.days

    # --- Page Routing ---
    if selected == "Home":
        render_command_center(portfolio_df, findings_df, team_df)
    elif selected == "Predictive Analytics":
        render_predictive_analytics(findings_df, portfolio_df)
    elif selected == "Systemic Risk":
        render_systemic_risk(findings_df, portfolio_df)
    elif selected == "PI Performance":
        render_pi_performance(portfolio_df, findings_df)
    elif selected == "Team & Strategy":
        render_organizational_capability(team_df, initiatives_df, audits_df, findings_df)
    elif selected == "Quality Impact":
        render_quality_impact(portfolio_df, findings_df)

if __name__ == "__main__":
    main()
